"""
Create PowerPoint Presentation for ASE Regression Analysis

Story Flow:
1. Section 1: Overall surrogate performance across descriptors
2. Section 2: PCA sensitivity analysis
3. Section 3: ORB descriptor deep dive (best PCA)

Emphasizes: R², Spearman, RMSE
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os
from pathlib import Path

def create_title_slide(prs, title, subtitle=""):
    """Create a title slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[0])  # Title slide layout

    title_shape = slide.shapes.title
    title_shape.text = title

    if subtitle and len(slide.placeholders) > 1:
        subtitle_shape = slide.placeholders[1]
        subtitle_shape.text = subtitle

    return slide

def create_section_slide(prs, section_title, description=""):
    """Create a section divider slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout

    # Add title
    left = Inches(0.5)
    top = Inches(2.5)
    width = Inches(9)
    height = Inches(1)

    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.text = section_title

    # Format title
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 51, 102)

    # Add description if provided
    if description:
        left = Inches(1)
        top = Inches(4)
        width = Inches(8)
        height = Inches(1)

        descBox = slide.shapes.add_textbox(left, top, width, height)
        tf = descBox.text_frame
        tf.text = description

        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(20)
        p.font.color.rgb = RGBColor(64, 64, 64)

    return slide

def create_content_slide(prs, title, image_path, notes=""):
    """Create a content slide with title and image."""
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout

    # Add title
    left = Inches(0.5)
    top = Inches(0.3)
    width = Inches(9)
    height = Inches(0.6)

    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.text = title

    p = tf.paragraphs[0]
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 51, 102)

    # Add image (centered, large)
    img_left = Inches(0.5)
    img_top = Inches(1.2)
    img_width = Inches(9)

    if os.path.exists(image_path):
        slide.shapes.add_picture(image_path, img_left, img_top, width=img_width)
    else:
        print(f"Warning: Image not found: {image_path}")

    # Add notes to slide notes section
    if notes:
        notes_slide = slide.notes_slide
        text_frame = notes_slide.notes_text_frame
        text_frame.text = notes

    return slide

def main():
    """Create the presentation."""

    print("=" * 80)
    print("CREATING PRESENTATION: ASE Regression Analysis")
    print("=" * 80)

    # Initialize presentation
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Get figures directory (use PNG versions)
    script_dir = Path(__file__).parent
    figures_dir = script_dir.parent / 'figures_png'

    # ========================================================================
    # TITLE SLIDE
    # ========================================================================
    print("\n1. Creating title slide...")
    create_title_slide(
        prs,
        "ASE Regression Analysis",
        "Surrogate Model Performance for Materials Property Prediction\n" +
        "n=500 Training Samples | 5-Fold Cross-Validation Holdout Performance"
    )

    # ========================================================================
    # SECTION 1: OVERALL SURROGATE PERFORMANCE
    # ========================================================================
    print("\n2. Creating Section 1: Overall Performance...")

    create_section_slide(
        prs,
        "Section 1: Surrogate Performance",
        "Comparing GP, MTGP, and DGP across MACE, ORB, SOAP, UMA descriptors"
    )

    # Heatmaps (averaged across properties, best PCA per surrogate)
    create_content_slide(
        prs,
        "Overall Performance: R² (Averaged Across Properties)",
        str(figures_dir / 'heatmaps' / 'averaged_R2_n500.png'),
        notes="Heatmap showing R² for each model-descriptor combination. " +
              "Uses best PCA per surrogate (GP/MTGP prefer PCA=50, DGP uses PCA=10-25). " +
              "ORB descriptor consistently best (R² ≈ 0.78-0.81)."
    )

    create_content_slide(
        prs,
        "Overall Performance: RMSE (Averaged Across Properties)",
        str(figures_dir / 'heatmaps' / 'averaged_RMSE_n500.png'),
        notes="RMSE (lower is better). ORB shows lowest error (~17-19), " +
              "SOAP highest error (~38-41)."
    )

    create_content_slide(
        prs,
        "Overall Performance: Spearman Correlation (Averaged)",
        str(figures_dir / 'heatmaps' / 'averaged_Spearman_n500.png'),
        notes="Spearman rank correlation (higher is better). " +
              "GP and MTGP show best ranking performance with ORB (0.84)."
    )

    # Bar charts (averaged)
    create_content_slide(
        prs,
        "R² Comparison Across Descriptors (Best PCA per Surrogate)",
        str(figures_dir / 'bar_charts' / 'averaged_R2_n500.png'),
        notes="Bar chart view of R². ORB descriptor clearly superior. " +
              "DGP slightly higher R² than GP with ORB (0.807 vs 0.803)."
    )

    create_content_slide(
        prs,
        "RMSE Comparison Across Descriptors",
        str(figures_dir / 'bar_charts' / 'averaged_RMSE_n500.png'),
        notes="RMSE comparison. GP+ORB shows lowest error (17.4), " +
              "followed closely by DGP+ORB (17.8)."
    )

    create_content_slide(
        prs,
        "Spearman Correlation Across Descriptors",
        str(figures_dir / 'bar_charts' / 'averaged_Spearman_n500.png'),
        notes="Spearman correlation. GP+ORB and MTGP+ORB both achieve 0.84, " +
              "outperforming DGP+ORB (0.83)."
    )

    # ========================================================================
    # SECTION 2: PCA SENSITIVITY ANALYSIS
    # ========================================================================
    print("\n3. Creating Section 2: PCA Sensitivity...")

    create_section_slide(
        prs,
        "Section 2: PCA Sensitivity",
        "How does PCA dimensionality affect surrogate performance?"
    )

    create_content_slide(
        prs,
        "PCA Sensitivity: R² (Averaged Across All Properties)",
        str(figures_dir / 'pca_sensitivity' / 'averaged_R2_n500.png'),
        notes="CRITICAL FINDING: DGP shows extreme PCA sensitivity! " +
              "With ORB: R² peaks at PCA=25 (0.81) but CRASHES to -0.02 at PCA=50. " +
              "GP/MTGP stable and improve with higher PCA. " +
              "This explains why DGP prefers lower PCA values."
    )

    create_content_slide(
        prs,
        "PCA Sensitivity: Spearman Correlation (Averaged)",
        str(figures_dir / 'pca_sensitivity' / 'averaged_Spearman_n500.png'),
        notes="Spearman correlation shows similar pattern. " +
              "DGP unstable at high PCA, GP/MTGP stable and monotonically improving."
    )

    create_content_slide(
        prs,
        "PCA Sensitivity Per Property: R² (ORB Descriptor)",
        str(figures_dir / 'pca_sensitivity' / 'per_property_R2_n500.png'),
        notes="Per-property PCA sensitivity with ORB descriptor. " +
              "DGP shows extreme sensitivity for ALL properties (ranges 0.6-0.95). " +
              "GP relatively stable (ranges < 0.12). " +
              "Hardest properties: poisson_ratio, elastic_anisotropy."
    )

    create_content_slide(
        prs,
        "PCA Sensitivity Per Property: Spearman (ORB Descriptor)",
        str(figures_dir / 'pca_sensitivity' / 'per_property_Spearman_n500.png'),
        notes="Spearman per-property sensitivity. " +
              "Same pattern: DGP highly sensitive, GP/MTGP stable."
    )

    # ========================================================================
    # SECTION 3: ORB DESCRIPTOR DEEP DIVE (BEST PCA)
    # ========================================================================
    print("\n4. Creating Section 3: ORB Deep Dive...")

    create_section_slide(
        prs,
        "Section 3: ORB Descriptor Analysis",
        "Deep dive into best-performing descriptor with optimal PCA per surrogate"
    )

    create_content_slide(
        prs,
        "Property Difficulty Matrix (ORB, Best PCA per Surrogate)",
        str(figures_dir / 'property_difficulty' / 'difficulty_matrix_per_surrogate_n500.png'),
        notes="3 heatmaps showing R² for each property with ORB descriptor. " +
              "GP uses PCA=50, MTGP uses PCA=50, DGP uses PCA=25. " +
              "EASY: Bulk moduli (K_Voigt, K_VRH, K_Reuss) - R² > 0.7 all models. " +
              "MODERATE: Shear moduli (G_*) - R² 0.5-0.7. " +
              "HARD: elastic_anisotropy, poisson_ratio - DGP struggles (R² 0.16-0.26)."
    )

    create_content_slide(
        prs,
        "Radar Chart: R² Across Properties (ORB, Best PCA)",
        str(figures_dir / 'radar_charts' / 'orb_R2_n500.png'),
        notes="Radar chart with properties as vertices, models as lines. " +
              "Uses best PCA per property: GP=PCA50 (all), MTGP=PCA50 (all), DGP=PCA25 (all). " +
              "Average R²: DGP=0.807 (highest), GP=0.803, MTGP=0.780. " +
              "All models best on K_Voigt, worst on elastic_anisotropy or poisson_ratio."
    )

    create_content_slide(
        prs,
        "Radar Chart: Spearman Across Properties (ORB, Best PCA)",
        str(figures_dir / 'radar_charts' / 'orb_Spearman_n500.png'),
        notes="Spearman radar chart. " +
              "Average Spearman: GP=0.844 (best), MTGP=0.843, DGP=0.826 (worst). " +
              "GP shows better ranking correlation despite slightly lower R² than DGP."
    )

    # Per-property bar charts (ORB only, show a few key properties)
    print("\n5. Adding per-property details (ORB only)...")

    key_properties = [
        ('K_Voigt', 'Easiest property to predict'),
        ('G_VRH', 'Moderate difficulty - shear modulus'),
        ('elastic_anisotropy', 'Hardest property - DGP struggles'),
        ('poisson_ratio', 'Very hard - derived property')
    ]

    for prop, description in key_properties:
        create_content_slide(
            prs,
            f"Property Deep Dive: {prop} ({description})",
            str(figures_dir / 'bar_charts' / 'per_property' / f'{prop}_R2_n500.png'),
            notes=f"R² for {prop} across all models and descriptors. " +
                  f"Shows best PCA per model-descriptor for THIS specific property. " +
                  f"{description}."
        )

    # ========================================================================
    # SUMMARY SLIDE
    # ========================================================================
    print("\n6. Creating summary slide...")

    slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title and content
    title = slide.shapes.title
    title.text = "Key Findings & Recommendations"

    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()

    findings = [
        "ORB descriptor consistently best across all models (R² ≈ 0.78-0.81)",
        "DGP shows extreme PCA sensitivity - CRASHES at PCA=50 (must use PCA≤25)",
        "GP provides best balance: high R² (0.803), best Spearman (0.844), stable",
        "Property difficulty: Bulk moduli (easy) > Shear moduli > Derived properties (hard)",
        "DGP struggles with poisson_ratio (R²=0.16) and elastic_anisotropy (R²=0.26)",
        "",
        "Recommendations:",
        "  • Use ORB descriptor for materials property prediction",
        "  • Use GP for reliable predictions with good uncertainty quantification",
        "  • Avoid DGP at high PCA (requires extensive hyperparameter tuning)",
        "  • PCA choice: GP/MTGP → PCA=50, DGP → PCA=25"
    ]

    for finding in findings:
        p = tf.add_paragraph()
        p.text = finding
        p.level = 0 if not finding.startswith('  ') else 1
        p.font.size = Pt(16)

    # ========================================================================
    # SAVE PRESENTATION
    # ========================================================================
    output_path = script_dir.parent / 'ASE_Regression_Analysis_Presentation.pptx'
    prs.save(str(output_path))

    print("\n" + "=" * 80)
    print(f"✓ Presentation saved to: {output_path}")
    print(f"  Total slides: {len(prs.slides)}")
    print("=" * 80)

    return output_path

if __name__ == '__main__':
    main()
